# pptx_copy_agent.py
from pptx import Presentation
from copy import deepcopy
from typing import List, Optional
from pptx.oxml.ns import qn
import shutil
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt
from typing import Iterable, Optional
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE



class PptxSlideCopyAgent:
    """
    Агент копіює вказаний слайд, опційно додає фонову картинку, змінює текст і нотатки.
    Використовує безпечний прийом: фон додається ПЕРШИМ shape-ом, порядок shapes не змінюється.
    """

    def __init__(
        self,
        pptx_path: str,
        out_path: str = "template_copy_with_bg.pptx",
        src_slide_index: int = 1,                 # 2-й слайд (індексація з 0)
        insert_after_src: bool = True,            # вставити копію одразу після джерела
    ):
        self.pptx_path = pptx_path
        self.out_path = out_path
        self.src_idx = src_slide_index
        self.insert_after_src = insert_after_src

        shutil.copy2(self.pptx_path, self.out_path)
        
    def autosize_text_boxes(self,
                            *,
                            path: str = None,
                            slides=None,
                            left_right_margin_pt: float = 2.0,
                            max_width=None,
                            save_to: str = None) -> str:
        """
        Після генерації презентації:
        - вимикає перенесення рядків (щоб мінялась ширина рамки),
        - вмикає SHAPE_TO_FIT_TEXT,
        - робить "поштовх" тексту (додає і одразу прибирає пробіл/символ), щоб PowerPoint перерахував розміри.
        Працює для шейпів, груп і таблиць.
        """
        from pptx import Presentation

        in_path = path or getattr(self, "out_path", None) or getattr(self, "pptx_path", None)
        if not in_path:
            raise ValueError("Не вказано шлях до презентації: передай path або задай self.out_path/self.pptx_path")

        prs = Presentation(in_path)

        idx_set = set(slides) if slides is not None else None

        def _bump_text_frame(tf):
            # мінімальні поля, щоб не заважали підгонці ширини
            try:
                tf.margin_left = Pt(left_right_margin_pt)
                tf.margin_right = Pt(left_right_margin_pt)
            except Exception:
                pass

            # критично: ширина під текст
            tf.word_wrap = False
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            # інколи корисно зафіксувати вертикальне вирівнювання (не обов'язково)
            try:
                tf.vertical_anchor = MSO_ANCHOR.TOP
            except Exception:
                pass

            # --- ПОШТОВХ ТЕКСТУ ---
            # Якщо просто виставити auto_size, PowerPoint часто “прокидається”
            # лише під час ручного редагування. Тож зімітуємо це:
            for p in tf.paragraphs:
                if p.runs:
                    r = p.runs[-1]
                else:
                    r = p.add_run()
                original = r.text
                # додаємо і прибираємо символ — цього достатньо, щоб Office перерахував розміри
                r.text = (original or "") + " "
                r.text = original

        def _walk_shapes(shapes):
            for sh in shapes:
                st = getattr(sh, "shape_type", None)
                # 1) звичайний текстовий шейп
                if getattr(sh, "has_text_frame", False):
                    tf = sh.text_frame
                    _bump_text_frame(tf)
                    if max_width is not None and sh.width > max_width:
                        sh.width = max_width

                # 2) група
                elif st == MSO_SHAPE_TYPE.GROUP:
                    _walk_shapes(sh.shapes)

                # 3) таблиця
                elif hasattr(sh, "has_table") and sh.has_table:
                    for row in sh.table.rows:
                        for cell in row.cells:
                            if cell.text_frame:
                                _bump_text_frame(cell.text_frame)

        for i, slide in enumerate(prs.slides):
            if idx_set is not None and i not in idx_set:
                continue
            _walk_shapes(slide.shapes)

        out_path = save_to or in_path
        prs.save(out_path)
        return out_path
    

    def _duplicate_slide_with_bg_first(self, prs: Presentation):
        src = prs.slides[self.src_idx]
        new_slide = prs.slides.add_slide(src.slide_layout)
    
        # 1) прибрати placeholders від layout (як у тебе)
        spTree = new_slide.shapes._spTree
        for sh in list(new_slide.shapes):
            spTree.remove(sh.element)
    
        # 2) фон першим (як у тебе)
        if self.bg_image_path:
            new_slide.shapes.add_picture(
                self.bg_image_path, 0, 0,
                width=prs.slide_width, height=prs.slide_height
            )
    
        # 3) скопіювати ВСІ шейпи (як у тебе)
        for sh in src.shapes:
            new_slide.shapes._spTree.insert_element_before(deepcopy(sh.element), 'p:extLst')
    
        # 3.1) ДОДАТКОВО: скопіювати transition та timing на рівні слайду
        sld_src = src._element
        sld_dst = new_slide._element
    
        # <p:transition> (якщо є)
        trans_src = sld_src.find(qn('p:transition'))
        if trans_src is not None:
            # видаляємо існуючий transition, якщо layout щось підсунув
            trans_dst = sld_dst.find(qn('p:transition'))
            if trans_dst is not None:
                sld_dst.remove(trans_dst)
            sld_dst.append(deepcopy(trans_src))
    
        # <p:timing> з усіма послідовностями анімацій
        timing_src = sld_src.find(qn('p:timing'))
        if timing_src is not None:
            timing_dst = sld_dst.find(qn('p:timing'))
            if timing_dst is not None:
                sld_dst.remove(timing_dst)
            sld_dst.append(deepcopy(timing_src))
    
        # 4) переставити новий слайд одразу після джерела (як у тебе)
        if self.insert_after_src:
            sldIdLst = prs.slides._sldIdLst
            sldIds = list(sldIdLst)
            new_el = sldIds[-1]
            sldIdLst.remove(new_el)
            sldIdLst.insert(self.src_idx + 1, new_el)
    
        return new_slide
    
    def _ensure_anim_updates_for_paragraphs(slide, shape):
        """
        Вмикає autoUpdateAnim для анімації 'By paragraph' на конкретному shape,
        щоб нові абзаци автоматично підхоплювали ту саму анімацію.
        """
        spid = str(shape.shape_id)  # саме за цим ID анімація прив’язується в <p:bldP spid="...">
        sld = slide._element
    
        timing = sld.find(qn('p:timing'))
        if timing is None:
            return
    
        bldLst = timing.find(qn('p:bldLst'))
        if bldLst is None:
            return
    
        # знайти build-параметри для нашого shape
        for bldP in bldLst.findall(qn('p:bldP')):
            if bldP.get('spid') == spid:
                # увімкнути автооновлення побудови
                bldP.set('autoUpdateAnim', '1')
                # якщо режим не вказаний — виставити "p" = by paragraph
                if bldP.get('build') is None:
                    bldP.set('build', 'p')
                # більше нічого не треба
                break


        
    def _apply_text_and_bullets(self, slide, prs):
        # 1) знайти перші два текстові шейпи (title + body)
        text_shapes = [sh for sh in slide.shapes if getattr(sh, "has_text_frame", False)]


        
        # ----------------------------------------------------------------------
        # Допоміжна функція – вимірює висоту текстового блоку
        # ----------------------------------------------------------------------
        
                
        def _set_para_text_keep_style(paragraph, new_text: str):
            if paragraph.runs:
                paragraph.runs[0].text = str(new_text)
                for r in paragraph.runs[1:]:
                    r.text = ""
            else:
                r = paragraph.add_run()
                r.text = str(new_text)
    
        def _clone_basic_font(from_run, to_run):
            fr, tr = from_run.font, to_run.font
            if fr.name:   tr.name = fr.name
            if fr.size:   tr.size = fr.size
            if fr.bold is not None:  tr.bold = fr.bold
            if fr.italic is not None: tr.italic = fr.italic
            try:
                if fr.color and getattr(fr.color, "rgb", None):
                    tr.color.rgb = fr.color.rgb
            except Exception:
                pass
                
        body_shape = None
        if len(text_shapes) >= 2:
            body_shape = text_shapes[1]
            
        # Title
        if len(text_shapes) >= 1:
            tf = text_shapes[0].text_frame
            if tf.paragraphs:
                _set_para_text_keep_style(tf.paragraphs[0], self.title_text)
    
        # Body bullets
        if len(text_shapes) >= 2 and self.bullets:
            tfb = text_shapes[1].text_frame
            existing = len(tfb.paragraphs)
            need = len(self.bullets)
    
            for i in range(min(existing, need)):
                p = tfb.paragraphs[i]
                _set_para_text_keep_style(p, self.bullets[i])
    
            if need > existing:
                ref_p = tfb.paragraphs[existing - 1] if existing else None
                ref_run0 = ref_p.runs[0] if (ref_p and ref_p.runs) else None
                for i in range(existing, need):
                    p_new = tfb.add_paragraph()
                    if ref_p is not None:
                        p_new.level = ref_p.level
                    r = p_new.add_run()
                    if ref_run0 is not None:
                        _clone_basic_font(ref_run0, r)
                    r.text = str(self.bullets[i])
    
                
    def _apply_notes(self, slide):
        if self.notes_text is None:
            return
        notes_tf = slide.notes_slide.notes_text_frame
        notes_tf.clear()
        for i, line in enumerate(str(self.notes_text).split("\n")):
            p = notes_tf.paragraphs[0] if i == 0 else notes_tf.add_paragraph()
            p.text = line
        
    # def run_intro(self,
    #     notes_text: Optional[str] = None
    #              ) -> str:  
    #     self.notes_text = notes_text
    #     prs = Presentation(self.out_path)
    #     slide = prs.slides[0]
    #     self._apply_notes(slide)
    #     prs.save(self.out_path)
    #     return self.out_path
    
    def run_intro(
        self,
        title_text: Optional[str] = None,
        body_text: Optional[str] = None,
        bg_image_path: Optional[str] = None,
        notes_text: Optional[str] = None,
        intro: Optional[str] = None
    ) -> str:
        """
        Оновлює перший слайд:
        • Фон — найнижчий шар
        • Текст у 1-му та 2-му блоках — ЗБЕРІГАЄ ВСІ СТИЛІ
        • Нотатки
        """
        self.bg_image_path = bg_image_path
        self.notes_text = notes_text
    
        prs = Presentation(self.out_path)
        if intro is None:
            slide = prs.slides[0]
        else:
            slide = prs.slides[-1]
    
        # === 1. ФОН НАЙНИЖЧИЙ ШАР ===
        if self.bg_image_path:
            spTree = slide.shapes._spTree
            # Видаляємо старі фони
            to_remove = []
            for child in spTree:
                if child.tag.endswith('}pic'):
                    xfrm = child.find('.//a:xfrm', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                    if xfrm is not None:
                        ext = xfrm.find('.//a:ext', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                        if ext is not None:
                            try:
                                cx = int(ext.get('cx') or 0)
                                cy = int(ext.get('cy') or 0)
                                if cx >= prs.slide_width * 0.9 and cy >= prs.slide_height * 0.9:
                                    to_remove.append(child)
                            except:
                                pass
            for el in to_remove:
                spTree.remove(el)
    
            picture = slide.shapes.add_picture(
                self.bg_image_path, 0, 0,
                width=prs.slide_width, height=prs.slide_height
            )
            spTree.insert(2, picture.element)  # ← найнижчий шар
    
        # === 2. ЗАМІНА ТЕКСТУ ЗІ ЗБЕРЕЖЕННЯМ СТИЛЮ ===
        text_shapes = [sh for sh in slide.shapes if getattr(sh, "has_text_frame", False)]
    
        def _set_text_preserve_style(paragraph, new_text: str):
            """Замінює текст, зберігаючи ВСІ стилі run-ів"""
            if not new_text:
                paragraph.text = ""
                return
    
            # Якщо run-ів немає — створюємо один
            if not paragraph.runs:
                run = paragraph.add_run()
                run.text = str(new_text)
                return
    
            # Беремо перший run як шаблон стилю
            template_run = paragraph.runs[0]
            template_run.text = str(new_text)
    
            # Очищаємо інші run-и, але не видаляємо (щоб не зламати форматування)
            for run in paragraph.runs[1:]:
                run.text = ""
    
        # --- Заголовок (1-й блок) ---
        if title_text is not None and len(text_shapes) >= 1:
            tf = text_shapes[0].text_frame
            if tf.paragraphs:
                _set_text_preserve_style(tf.paragraphs[0], title_text)
    
        # --- Тіло (2-й блок) ---
        if body_text is not None and len(text_shapes) >= 2:
            tf = text_shapes[1].text_frame
            lines = [line.strip() if line.strip() else " " for line in str(body_text).split('\n')]
            existing = tf.paragraphs
            need = len(lines)
    
            # 1. Замінюємо існуючі абзаци
            for i in range(min(len(existing), need)):
                _set_text_preserve_style(existing[i], lines[i])
    
            # 2. Додаємо нові абзаци (копіюємо стиль з останнього)
            if need > len(existing):
                ref_para = existing[-1] if existing else None
                for i in range(len(existing), need):
                    new_para = tf.add_paragraph()
                    if ref_para is not None:
                        # Копіюємо рівень, вирівнювання, маркер
                        new_para.level = ref_para.level
                        new_para.alignment = ref_para.alignment
                        # Копіюємо стиль шрифту з першого run
                        if ref_para.runs:
                            ref_run = ref_para.runs[0]
                            if new_para.runs:
                                target_run = new_para.runs[0]
                            else:
                                target_run = new_para.add_run()
                            # --- Копіюємо ВСІ властивості шрифту ---
                            font = ref_run.font
                            target = target_run.font
                            if font.name:      target.name = font.name
                            if font.size:      target.size = font.size
                            if font.bold is not None:   target.bold = font.bold
                            if font.italic is not None: target.italic = font.italic
                            if font.underline is not None: target.underline = font.underline
                            if getattr(font.color, 'rgb', None) is not None:
                                target.color.rgb = font.color.rgb
                            if font.fill and hasattr(font.fill, 'solid'):
                                try:
                                    target.fill.solid()
                                    target.fill.fore_color.rgb = font.fill.fore_color.rgb
                                except:
                                    pass
                    _set_text_preserve_style(new_para, lines[i])
    
            # 3. Видаляємо зайві абзаци (з кінця)
            while len(tf.paragraphs) > need:
                tf.paragraphs[-1]._element.getparent().remove(tf.paragraphs[-1]._element)
    
        # === 3. НОТАТКИ ===
        if notes_text is not None:
            self._apply_notes(slide)
    
        prs.save(self.out_path)
        return self.out_path
    
    
    def delete_slide(self, index: int) -> None:
        """Видаляє слайд за індексом (0-based)."""
        prs = Presentation(self.out_path)

        # Дістаємо список ідентифікаторів слайдів
        slide_id_list = prs.slides._sldIdLst
        slide_id = slide_id_list[index]
        rId = slide_id.rId  # id зв’язку на slide-part

        # Видаляємо відношення і сам слайд з колекції
        prs.part.drop_rel(rId)
        slide_id_list.remove(slide_id)

        prs.save(self.out_path)
        return self.out_path
    
    def run(self,
        bg_image_path: Optional[str] = None,      # шлях до зображення або None
        title_text: str = "Test 1",
        bullets: Optional[List[str]] = None,
        notes_text: Optional[str] = None,
           ) -> str:
        self.bg_image_path = bg_image_path
        self.title_text = title_text
        self.bullets = bullets or ["пункт 1", "пункт 2", "пункт 3"]
        self.notes_text = notes_text
        
        prs = Presentation(self.out_path)
        slide = self._duplicate_slide_with_bg_first(prs)
        self._apply_text_and_bullets(slide, prs)
        self._apply_notes(slide)
        prs.save(self.out_path)
        return self.out_path


if __name__ == "__main__":
    # приклад використання
    agent = PptxSlideCopyAgent(
        pptx_path="template.pptx",
        out_path="template_copy2_with_bg.pptx",
        src_slide_index=1,                  # копіюємо 2-й слайд
        bg_image_path="image.png",          # фон (None, якщо не потрібно)
        title_text="Test 1",
        bullets=["пункт 1", "пункт 2", "пункт 3"],
        notes_text="Це приклад тексту в нотатках до цього слайду.",
        insert_after_src=True
    )
    result = agent.run()
    print(f"Збережено у {result}")
